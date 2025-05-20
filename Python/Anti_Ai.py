import tkinter as tk
from tkinter import filedialog, messagebox
import os
import filetype
from PIL import Image
from pptx import Presentation
from pdfminer.high_level import extract_text
from tkinter import ttk

def detect_ai(file_path):
    kind = filetype.guess(file_path)
    if kind:
        file_type = kind.mime
    else:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            file_type = "application/pdf"
        elif ext in [".ppt", ".pptx"]:
            file_type = "application/vnd.ms-powerpoint"
        elif ext in [".jpg", ".jpeg", ".png"]:
            file_type = "image"
        elif ext in [".mp4", ".mov", ".avi"]:
            file_type = "video"
        else:
            return "Unsupported format"

    # Mock detection logic
    result = ""
    if file_type.startswith("image"):
        result = detect_ai_image(file_path)
    elif file_type.startswith("application/pdf"):
        result = detect_ai_pdf(file_path)
    elif "powerpoint" in file_type:
        result = detect_ai_ppt(file_path)
    elif file_type.startswith("video"):
        result = detect_ai_video(file_path)
    else:
        result = "Detection not available."

    return result

def detect_ai_image(file_path):
    try:
        with Image.open(file_path) as img:
            metadata = img.info
            if "software" in metadata and "AI" in metadata["software"]:
                return "Likely AI-generated image"
            # else:

            #     for filename in os.listdir(file_path):
            #         if filename.endswith('.jpg') or filename.endswith('.png'):
            #             print(filename)

            #     # if "AI" in title or "ChatGPT" in title or "generated" in title:
            #     #     return "Likely AI-generated image"
            return "Likely human-made image"
    except Exception:
        return "Error reading image"

def detect_ai_pdf(file_path):
    try:
        text = extract_text(file_path)
        if "ChatGPT" in text or "AI generated" in text:
            return "Likely AI-generated PDF"
        return "Likely human-written PDF"
    except Exception:
        return "Could not read PDF"

def detect_ai_ppt(file_path):
    try:
        prs = Presentation(file_path)
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        if "ChatGPT" in all_text or "AI" in all_text:
            return "Likely AI-generated PPT"
        return "Likely human-made PPT"
    except Exception:
        return "Could not read PPT"

def detect_ai_video(file_path):
    return "Video analysis not implemented. Needs AI model."

def show_tooltip(widget, text):
    tooltip = tk.Toplevel(widget)
    tooltip.wm_overrideredirect(True)
    tooltip.wm_geometry(f"+{widget.winfo_rootx()+20}+{widget.winfo_rooty()+20}")
    label = tk.Label(tooltip, text=text, background="#333", foreground="#fff", borderwidth=1, relief="solid", font=("Helvetica", 9))
    label.pack()
    widget.tooltip = tooltip

def hide_tooltip(widget):
    if hasattr(widget, 'tooltip'):
        widget.tooltip.destroy()
        del widget.tooltip

def upload_file():

    file_path = filedialog.askopenfilename(
        title="Select a file",
        filetypes=[("All Supported", "*.pdf *.ppt *.pptx *.jpg *.jpeg *.png *.mp4 *.mov *.avi"),
                   ("PDF files", "*.pdf"),
                   ("PowerPoint files", "*.ppt *.pptx"),
                   ("Image files", "*.jpg *.jpeg *.png"),
                   ("Video files", "*.mp4 *.mov *.avi"),
                   ("All files", "*.*")]
    
    )
    if file_path:
        progress_bar.start()
        root.update_idletasks()
        result = detect_ai(file_path)
        progress_bar.stop()
        result_label.config(text=f"Detection Result:\n{result}")
    return file_path

# Tkinter GUI setup
root = tk.Tk()
root.title("AI Content Detector")
root.geometry("520x400")
root.config(bg="#eaf6fb")
root.resizable(False, False)

# Apply a theme
style = ttk.Style()
style.theme_use('clam')
style.configure('TButton', font=('Helvetica', 12, 'bold'), foreground='#fff', background='#007acc', padding=10)
style.map('TButton', background=[('active', '#005f99')])
style.configure('TLabel', font=('Helvetica', 12), background='#eaf6fb')
style.configure('Title.TLabel', font=('Helvetica', 18, 'bold'), background='#eaf6fb', foreground='#007acc')
style.configure('Footer.TLabel', font=('Helvetica', 10), background='#eaf6fb', foreground='#888')
style.configure('TFrame', background='#eaf6fb')
style.configure('HighlightFooter.TLabel', font=('Helvetica', 13, 'bold'), background='#eaf6fb', foreground='#d9534f')

# Title
title = ttk.Label(root, text="AI vs Real Content Detector", style='Title.TLabel')
title.pack(pady=(25, 5))

# Supported file formats info
formats_info = ttk.Label(
    root,
    text="Supported formats: PDF, PowerPoint (.ppt, .pptx), Images (.jpg, .jpeg, .png), Videos (.mp4, .mov, .avi)",
    style='TLabel',
    wraplength=480,
    anchor='center',
    justify='center'
)
formats_info.pack(pady=(0, 15))

# Upload button
upload_btn = ttk.Button(root, text="Upload File", command=upload_file)
upload_btn.pack(pady=15)

# Result label
result_label = ttk.Label(root, text="", style='TLabel', wraplength=400, anchor='center', justify='center')
result_label.pack(pady=25)

# Footer
footer = ttk.Label(root, text="© 2025 AI Detector", style='Footer.TLabel')
footer.pack(side='bottom', pady=(5, 0))

# Improvement note at the very bottom
improve_note = ttk.Label(
    root,
    text="We are working on it for further improvement",
    style='Footer.TLabel',
    anchor='center',
    justify='center'
)
improve_note.pack(side='bottom', pady=(0, 10))

# Progress bar (hidden by default)
progress_bar = ttk.Progressbar(root, mode='indeterminate')
progress_bar.pack(fill='x', padx=10, pady=10)

# Start the Tkinter event loop
root.mainloop()